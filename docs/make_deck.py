"""Generate the consolidated talk deck: docs/RWE_talk.pptx.

Recreates Bibek Paudel's WWW'21 talk as a clean, editable slide deck *and*
foregrounds how each part maps to this repository (a "-> In the code" note per
slide), embedding the project's own diagrams from ``docs/images/``.

This is a recreation from the talk's content -- not a copy of the original
slides (their image files are not available). Run with::

    python docs/make_deck.py

Requires ``python-pptx`` (the ``[docs]`` extra). Deterministic: re-running
reproduces an identical deck.
"""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
IMAGES = HERE / "images"
OUT = HERE / "RWE_talk.pptx"

# Palette (matches the talk's maroon + the diagrams).
MAROON = RGBColor(0x8B, 0x1A, 0x1A)
RULE = RGBColor(0xCC, 0xCC, 0xCC)
INK = RGBColor(0x33, 0x33, 0x33)
MUTE = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x55, 0xA8, 0x68)
GREEN_DK = RGBColor(0x2E, 0x6B, 0x3E)
PANEL = RGBColor(0xEE, 0xF5, 0xEE)

SW, SH = 13.333, 7.5  # 16:9 inches


# --------------------------------------------------------------------------- #
# Slide-building helpers
# --------------------------------------------------------------------------- #
def _title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(SW - 1.2), Inches(0.95))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = MAROON
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.28),
                                  Inches(SW - 1.2), Pt(2))
    rule.fill.solid()
    rule.fill.fore_color.rgb = RULE
    rule.line.fill.background()
    rule.shadow.inherit = False


def _bullets(slide, bullets, top, height, width=SW - 1.2, left=0.6, size=19):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        prefix = "•  " if level == 0 else "       –  "
        r = p.add_run()
        r.text = prefix + text
        r.font.size = Pt(size - 2 * level)
        r.font.color.rgb = INK if level == 0 else MUTE


def _image_fit(slide, name, box_top, box_bottom, max_width=11.0):
    """Place an image, scaled to fit (and centred in) the given vertical box."""
    w_px, h_px = Image.open(IMAGES / name).size
    aspect = w_px / h_px
    max_h = box_bottom - box_top
    width = min(max_width, max_h * aspect)
    height = width / aspect
    left = (SW - width) / 2
    top = box_top + (max_h - height) / 2
    slide.shapes.add_picture(str(IMAGES / name), Inches(left), Inches(top),
                             width=Inches(width))


def _code_note(slide, note):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6),
                                 Inches(6.55), Inches(SW - 1.2), Inches(0.62))
    shp.fill.solid()
    shp.fill.fore_color.rgb = PANEL
    shp.line.color.rgb = GREEN
    shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "→ In the code:  "
    r1.font.size = Pt(13)
    r1.font.bold = True
    r1.font.color.rgb = GREEN_DK
    r2 = p.add_run()
    r2.text = note
    r2.font.size = Pt(13)
    r2.font.color.rgb = INK


def slide(prs, title, bullets, image=None, image_box=(3.35, 6.4),
          code_note=None, notes=None, bullets_top=1.55, bullets_height=4.6):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title(s, title)
    _bullets(s, bullets, bullets_top, bullets_height)
    if image:
        _image_fit(s, image, image_box[0], image_box[1])
    if code_note:
        _code_note(s, code_note)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    box = s.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(SW - 1.8), Inches(1.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Random Walks with Erasure"
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.color.rgb = MAROON
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Diversifying Personalized Ranking"
    r2.font.size = Pt(26)
    r2.font.color.rgb = INK
    sub = s.shapes.add_textbox(Inches(0.9), Inches(4.2), Inches(SW - 1.8), Inches(2.0))
    tf2 = sub.text_frame
    tf2.word_wrap = True
    for txt, sz, col in [
        ("Paper: Bibek Paudel (Stanford) · Abraham Bernstein (UZH) — WWW 2021", 18, MUTE),
        ("A slide-by-slide walkthrough mapped to this implementation.", 16, MUTE),
        ("Deck recreated from the talk; diagrams + code-mapping by the project.", 13, RULE),
    ]:
        p = tf2.paragraphs[0] if txt.startswith("Paper") else tf2.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.color.rgb = col
    s.notes_slide.notes_text_frame.text = (
        "Implementation of Paudel & Bernstein, WWW 2021. This deck recreates the "
        "talk and maps every part to the rwe/ codebase.")
    return s


# --------------------------------------------------------------------------- #
# The deck
# --------------------------------------------------------------------------- #
def build():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    title_slide(prs)

    slide(prs, "The problem: filter bubbles", [
        ("Recommenders show \"more of what you already like\" → filter bubbles.", 0),
        ("2010–13: social media praised as an enabler of democracy and protest.", 0),
        ("Then criticised for increasing polarization and harming society.", 0),
        ("High polarization → suspicion, hostility, and lower political participation.", 0),
    ], code_note="motivation captured in GUIDE.md §1",
        notes="Ten years ago social media was praised as a democracy enabler; "
              "soon after, blamed for polarization. We focus on the diversity of "
              "information in personalized recommendations.")

    slide(prs, "How to bridge the divide?", [
        ("Cross-cutting awareness = exposure to viewpoints from the other side.", 0),
        ("It is linked to greater political understanding and participation.", 0),
        ("Goal: recommend content that is relevant AND bridges viewpoints.", 0),
    ], code_note="RWEB (bridging) + extensions satisfaction.py / agent_sim.py",
        notes="Cross-cutting awareness and discussions increase political "
              "understanding and participation.")

    slide(prs, "Challenge of political content", [
        ("Naive fix: mix outlets by their known slant (e.g. AllSides). Not enough.", 0),
        ("Two articles from the SAME outlet were shared by opposite Brexit camps.", 0),
        ("Ideological positions are issue/event-specific → must be learned, not labelled.", 0),
    ], code_note="IdeologyModel.fit(R, S) learns issue-specific positions from behaviour",
        notes="Pre-existing labels / known slants are unreliable; positions "
              "change with the issue and event, so we learn them.")

    slide(prs, "Our approach — three steps", [
        ("1. Identify ideological positions of elites AND content from event discussions.", 0),
        ("2. Choose a diversification strategy (naive \"mix both sides\" may fail).", 0),
        ("3. Random Walk with Erasure (RWE): recommend diverse content via weak ties.", 0),
    ], code_note="ideology.py  ·  RWED / RWEB  ·  RWE",
        notes="Three parts: learn positions, pick a strategy, run RWE.")

    slide(prs, "Ideology of users, elites & content", [
        ("Each retweet yields three roles: User (retweeter), Elite (author), Content (URL).", 0),
        ("Two graphs: elite-endorsement R and content-share S.", 0),
        ("Place users, elites and content on one shared left–right scale.", 0),
    ], image="ideology_scale.png", bullets_height=1.4,
        code_note="IdeologyModel.fit(R, S) → positions θ (users), φ (elites), ψ (content)",
        notes="A retweet gives a user→elite endorsement (R) and a user→content "
              "endorsement (S); both place everyone on one 1-D scale.")

    slide(prs, "Ideal-point model + joint learning", [
        ("Elite:    p(R=1) = σ( −‖θ_u − φ_e‖² + α_u + β_e )", 0),
        ("Content:  p(S=1) = σ( −‖θ_u − ψ_c‖² + α_u + γ_c )", 0),
        ("Closer ideology ⇒ more likely to endorse.", 1),
        ("Joint objective: maximise  μ·Σ_R[·] + Σ_S[·]  −  (λ/2)(‖θ‖+‖φ‖+‖ψ‖)", 0),
        ("⟹ sim(u, c): similarity of political stance.", 0),
    ], code_note="ideology.py: Pi_R / Pi_S (eqs 6/9), _objective (eq 11); sim = RWEB.similarity",
        notes="Endorsement probability decreases with squared ideological "
              "distance; elite and content models are learned jointly.")

    slide(prs, "Estimated positions (results)", [
        ("Recovers left/right correctly across UK, US and Germany.", 0),
        ("Captures issue-specific cases (e.g. Conservatives who backed Remain → left).", 0),
        ("Joint learning (elite + content) beats elite-endorsement alone.", 0),
    ], code_note="validated on synthetic data (real Twitter data not shareable) — demo_synthetic.py",
        notes="Issue-specific positioning is exactly what learning-from-behaviour "
              "captures, unlike fixed party labels.")

    slide(prs, "Normative goals for diversification", [
        ("Personalized: high accuracy.", 0),
        ("Long-tail: expose users to more low-degree items.", 0),
        ("Bridging: show the other side — but not too polarizing — via weak-tie \"bridges\".", 0),
    ], code_note="accuracy metrics · RWED · RWEB (max_distance = \"not too far\")",
        notes="Diversity is calibrated: accurate, strategy-specific, and reachable "
              "via weak ties rather than jarring opposite extremes.")

    slide(prs, "Diversification strategies — the erase matrix Q", [
        ("Q on the user-item graph prefers certain random walks over others.", 0),
        ("Example I — Bridging: items on the opposite side of the user.", 1),
        ("Example II — Long-tail: low-degree items.", 1),
        ("Example III — Contrarian: items opposite to the user.", 1),
        ("Generalized framework: any property defined on items or user-item pairs.", 0),
    ], code_note="RWE(item_erasure=…): RWEB · RWED · or any callable (contrarian = one-liner)",
        notes="Q is the single knob; different Q gives different strategies. The "
              "framework is open-ended.")

    slide(prs, "Random Walk with Erasure — the algorithm", [
        ("Bipartite graph: m users, n items; k-step transitions Pᵏ.", 0),
        ("Erase matrix Q ∈ [0,1)^{(m+n)×(m+n)}.", 0),
        ("Sample multiple rounds; erased mass becomes the next round's start.", 0),
    ], image="bipartite_graph.png", bullets_height=1.4,
        code_note="FeedbackGraph (A^G, P=D⁻¹A^G); item_distribution(users, k)",
        notes="Normal k-step walk via the transition matrix, plus an erase matrix "
              "Q; at each round the initial-state probabilities are modified by Q.")

    slide(prs, "How erasure steers recommendations", [
        ("Score(item) = p·(1−q) / (1 − Σ p·q)   —   the mass that survives the \"tax\".", 0),
        ("High tax q ⇒ suppressed;  low tax q ⇒ recommended.", 0),
        ("Final score:  w_s = Σᵢ v_{s,i} Pᵏ ∘ (1 − Q).", 0),
    ], image="rwe_flow.png", bullets_height=1.4,
        code_note="RWE.score_iterative (rounds) == derived closed form; tested to match",
        notes="The repeat-forever erasure loop telescopes into a closed form; the "
              "denominator is constant per user so it doesn't change the ranking.")

    slide(prs, "Long-tail worked example (reproduced)", [
        ("Pᵏ_{u₄,i*} = [0.16, 0.16, 0.66];   Q = 1 − [1/3, 1/2, 1/2] = [0.66, 0.5, 0.5].", 0),
        ("Erased = Pᵏ ∘ Q;  start mass falls  1 → 0.51 → 0.26  across rounds.", 0),
        ("Suppresses the high-degree item, lifts the low-degree ones.", 0),
        ("Our code reproduces these numbers to the decimal.", 0),
    ], code_note="RWED (β=1) via score_iterative — verified against the slide",
        notes="A plain walk scores i1 and i2 equally; degree-based erasure lowers "
              "the high-degree node and raises low-degree ones.")

    slide(prs, "Experimental evaluation — three questions", [
        ("Q1. Does joint learning identify ideological positions accurately?", 0),
        ("Q2. Can RWE generate accurate AND diverse long-tail recommendations?", 0),
        ("Q3. Can RWE generate accurate AND ideologically diverse recommendations?", 0),
    ], code_note="IdeologyModel · RWED · RWEB  +  metrics.py",
        notes="Three evaluation axes: ideology accuracy, long-tail diversity, "
              "ideological diversity.")

    slide(prs, "Datasets & baselines", [
        ("6 Twitter graphs: UK / US / DE × {retweets (RT), URL shares}.", 0),
        ("Plus benchmarks ML-1M and Yelp for long-tail diversity.", 0),
        ("Baselines (RecSys'19 study): Collaborative Filtering, Matrix Factorization,", 0),
        ("graph-based P³, and long-tail-diversifying RP³_β.", 1),
    ], code_note="ItemKNN · BPRMF · P3 · RP3Beta;  data.py (synthetic stand-ins for private data)",
        notes="State-of-the-art recommenders from a recent evaluation study.")

    slide(prs, "Result I — ideological range of top-10", [
        ("Average spread of ideological positions in the top-10 list.", 0),
        ("RWE-B has the widest range on all but the smallest dataset.", 0),
        ("→ more viewpoint-spread per recommendation list.", 1),
    ], code_note="metrics.rec_range_at_k  (reported as rec_range@10)",
        notes="RWE-B's recommendations span a wider ideological range than the most "
              "competitive baselines.")

    slide(prs, "Result II — ideological diversity", [
        ("x = user position, y = mean recommended position (each dot a user).", 0),
        ("Baselines keep users on their OWN side (dots on the diagonal).", 0),
        ("RWE moves recommendations to the OPPOSITE quadrant / centre.", 0),
    ], image="quadrant_scatter.png", bullets_height=1.4,
        code_note="docs/make_diagrams.py quadrant_scatter — metrics.mean_recommended_position",
        notes="Baselines recommend from the same quadrant; RWE recommends from the "
              "opposite quadrants and the centre.")

    slide(prs, "Result III — ideological shift", [
        ("User-Shift = Pos(Recs) − Pos(u);   Train-Shift = Pos(Recs) − Pos(Train).", 0),
        ("Left users shift right, right users shift left.", 0),
        ("RWE shows the strongest, cleanest shift toward the opposite side.", 0),
    ], code_note="metrics.ideological_shift (exact definition) → directed_shift",
        notes="Shift = difference between recommended position and the user's own / "
              "training position.")

    slide(prs, "Result IV — weighted ideological diversity", [
        ("Two properties: lean toward centre (per user); wider range for extremes.", 0),
        ("Measures UW / TW (weighted by user / training position): Recs, Shift, Range.", 0),
        ("RWE-B is best on all three — while keeping accuracy (AUC, HR@10) on par.", 0),
    ], code_note="weighted_position (UW/TW-Recs) · weighted_shift · weighted_range — uw_* in evaluate()",
        notes="Lower is better except for UW-Range, AUC and HR@10. RWE-B at "
              "erasure exponent v ∈ {5, 2, 0.5}.")

    slide(prs, "Contributions & main results", [
        ("A probabilistic model to estimate ideological positions (users, elites, content).", 0),
        ("Personalized ranking based on RWE.", 0),
        ("Two diversification strategies: long-tail and political bridging.", 0),
        ("A general framework to diversify personalized ranking.", 0),
    ], code_note="ideology.py · RWE · RWED/RWEB · RWE(item_erasure=…)",
        notes="The four contributions, each realised in the rwe/ package.")

    slide(prs, "This implementation", [
        ("Full paper in rwe/ (graph, RWE, ideology, baselines, metrics) + 2 extensions.", 0),
        ("63 tests — including the worked erasure example reproduced to the decimal.", 0),
        ("5 diagrams, a beginner guide (GUIDE.md), and a slide-by-slide map (docs/TALK.md).", 0),
        ("Run:  pip install -e .   ·   pytest -q   ·   python examples/demo_synthetic.py", 0),
    ], code_note="github: greenwichg/random_walks_with_erasure",
        notes="What this repository adds on top of the paper.")

    slide(prs, "Thanks", [
        ("Paper: Paudel & Bernstein, \"Random Walks with Erasure\", WWW 2021.", 0),
        ("Original talk: Bibek Paudel — @biasedcoin · bibekp@stanford.edu.", 0),
        ("This deck, diagrams and code-mapping: the implementation project.", 0),
    ], notes="Acknowledging the Hasler Foundation's support of the original work.")

    prs.save(str(OUT))
    print(f"wrote {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
