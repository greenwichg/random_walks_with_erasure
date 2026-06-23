"""Render docs/RWE_talk.pptx to docs/RWE_talk.pdf (one page per slide).

A lightweight, dependency-light fallback PDF exporter: it reads the deck's
shapes (text boxes, pictures, filled rectangles) directly from the ``.pptx`` and
draws each slide to a PDF page with matplotlib. Use it to keep the shared PDF in
sync with the deck without LibreOffice / PowerPoint.

Run with::

    python docs/make_pdf.py

For a pixel-perfect export, open ``RWE_talk.pptx`` in PowerPoint / Keynote /
Google Slides and "Save as PDF"; this script gives a faithful, reproducible
match (all deck bullets fit on one line, so there is no text-wrapping drift).
"""

import io
from pathlib import Path

from pptx import Presentation
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.patches import FancyBboxPatch, Rectangle
from PIL import Image

HERE = Path(__file__).resolve().parent
PPTX = HERE / "RWE_talk.pptx"
PDF = HERE / "RWE_talk.pdf"
EMU = 914400.0  # English Metric Units per inch


def _rgb(color):
    try:
        return "#%02x%02x%02x" % (color[0], color[1], color[2])
    except Exception:
        return None


def _draw_slide(slide, sw, sh):
    fig = plt.figure(figsize=(sw, sh))
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, sw)
    ax.set_ylim(sh, 0)  # y inverted so (0,0) is top-left like a slide
    ax.axis("off")
    ax.add_patch(Rectangle((0, 0), sw, sh, fc="white", ec="#dddddd", lw=1))
    for sh_ in slide.shapes:
        left, top = sh_.left / EMU, sh_.top / EMU
        w, h = (sh_.width or 0) / EMU, (sh_.height or 0) / EMU
        if sh_.shape_type == 13:  # picture
            try:
                img = Image.open(io.BytesIO(sh_.image.blob))
                ax.imshow(img, extent=[left, left + w, top + h, top],
                          aspect="auto", zorder=2)
            except Exception:
                pass
            continue
        try:  # filled auto-shapes (the rule line, the code-note box)
            if sh_.fill.type == 1:
                ax.add_patch(FancyBboxPatch((left, top), w, h, boxstyle="round,pad=0.02",
                                            fc=_rgb(sh_.fill.fore_color.rgb) or "none",
                                            ec="none", zorder=1))
        except Exception:
            pass
        if sh_.has_text_frame and sh_.text_frame.text.strip():
            y = top + 0.18
            for p in sh_.text_frame.paragraphs:
                text = "".join(r.text for r in p.runs)
                if not text:
                    continue
                run = p.runs[0]
                size = run.font.size.pt if run.font.size else 18
                color = (_rgb(run.font.color.rgb)
                         if run.font.color and run.font.color.type is not None
                         else "#333333")
                ax.text(left + 0.1, y, text, fontsize=size, color=color or "#333333",
                        weight="bold" if run.font.bold else "normal", va="top", zorder=3)
                y += size / 72 * 1.45
    return fig


def main():
    prs = Presentation(str(PPTX))
    sw, sh = prs.slide_width / EMU, prs.slide_height / EMU
    n = 0
    with PdfPages(str(PDF)) as pdf:
        for slide in prs.slides:
            fig = _draw_slide(slide, sw, sh)
            pdf.savefig(fig)
            plt.close(fig)
            n += 1
    print(f"wrote {PDF}  ({n} pages)")


if __name__ == "__main__":
    main()
